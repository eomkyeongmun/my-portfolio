"""
슬라이드 빌더 헬퍼
- 텍스트 박스, 카드, 태그, 구분선 등 재사용 컴포넌트
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

from theme import (
    SLIDE_W, SLIDE_H,
    COLOR_BG, COLOR_FG, COLOR_MUTED, COLOR_ACCENT, COLOR_ACCENT_SOFT,
    COLOR_LINE, COLOR_CARD,
    SIZE_TITLE, SIZE_SUBTITLE, SIZE_H2, SIZE_H3, SIZE_BODY, SIZE_SMALL, SIZE_TAG,
    font_for,
)


def new_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs: Presentation):
    blank = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(blank)
    paint_background(slide, COLOR_BG)
    return slide


def paint_background(slide, color: RGBColor):
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, SLIDE_W, SLIDE_H)
    bg.line.fill.background()
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    # 배경은 항상 맨 뒤
    spTree = bg._element.getparent()
    spTree.remove(bg._element)
    spTree.insert(2, bg._element)
    return bg


def add_text(
    slide,
    text: str,
    left, top, width, height,
    *,
    font: str = "Pretendard",
    size: Pt = SIZE_BODY,
    color: RGBColor = COLOR_FG,
    bold: bool = False,
    align: PP_ALIGN = PP_ALIGN.LEFT,
    anchor: MSO_ANCHOR = MSO_ANCHOR.TOP,
    line_spacing: float = 1.3,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Emu(0)
    tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor

    lines = text.split("\n") if text else [""]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.line_spacing = line_spacing
        run = p.add_run()
        run.text = line
        run.font.name = font
        run.font.size = size
        run.font.bold = bold
        run.font.color.rgb = color
    return box


def add_rect(
    slide, left, top, width, height,
    *,
    fill: RGBColor = COLOR_CARD,
    line: RGBColor | None = None,
    line_width: float = 0.75,
    radius: float = 0.04,  # 0~0.5 (도형 모서리 라운드 비율)
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    # adjust rounded corner radius
    try:
        shape.adjustments[0] = radius
    except Exception:
        pass
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(line_width)
    return shape


def add_line(slide, left, top, width, *, color=COLOR_LINE, weight: float = 0.75):
    line = slide.shapes.add_connector(1, left, top, left + width, top)
    line.line.color.rgb = color
    line.line.width = Pt(weight)
    return line


def add_tag(slide, text: str, left, top, *, font: str = "Pretendard"):
    """소형 태그 칩 (자동 너비)"""
    # 자동 너비: 글자 수 기반 추정
    width = Inches(0.10 + len(text) * 0.085)
    height = Inches(0.28)
    chip = add_rect(slide, left, top, width, height,
                    fill=COLOR_ACCENT_SOFT, line=None, radius=0.5)
    add_text(slide, text, left, top, width, height,
             font=font, size=SIZE_TAG, color=COLOR_ACCENT, bold=True,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return left + width  # 다음 태그가 시작될 x 위치 반환


def add_header(slide, title: str, subtitle: str | None, lang: str):
    """슬라이드 상단 헤더 (제목 + 좌측 강조바)"""
    font = font_for(lang)
    # 좌측 강조 바
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                 Inches(0.6), Inches(0.55),
                                 Inches(0.08), Inches(0.55))
    bar.line.fill.background()
    bar.fill.solid()
    bar.fill.fore_color.rgb = COLOR_ACCENT
    # 제목
    add_text(slide, title,
             Inches(0.85), Inches(0.45),
             Inches(11.5), Inches(0.7),
             font=font, size=SIZE_TITLE, color=COLOR_FG, bold=True)
    # 서브타이틀
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.85), Inches(1.05),
                 Inches(11.5), Inches(0.4),
                 font=font, size=SIZE_SUBTITLE, color=COLOR_MUTED)
    # 하단 구분선
    add_line(slide, Inches(0.6), Inches(1.55), Inches(12.13))


def add_footer(slide, text: str, lang: str):
    add_text(slide, text,
             Inches(0.6), Inches(7.15),
             Inches(12.13), Inches(0.25),
             font=font_for(lang), size=SIZE_SMALL, color=COLOR_MUTED,
             align=PP_ALIGN.RIGHT)


def add_image_safe(slide, image_path: str, left, top, width, height):
    """이미지가 존재할 때만 삽입, 없으면 회색 박스로 대체"""
    import os
    if image_path and os.path.exists(image_path):
        return slide.shapes.add_picture(image_path, left, top, width=width, height=height)
    placeholder = add_rect(slide, left, top, width, height,
                           fill=RGBColor(0xF3, 0xF4, 0xF6), line=COLOR_LINE, radius=0.02)
    add_text(slide, "(image)", left, top, width, height,
             size=SIZE_SMALL, color=COLOR_MUTED,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return placeholder
